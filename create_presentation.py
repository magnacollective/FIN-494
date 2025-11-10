from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# Create presentation
prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(7.5)

def add_title_slide(prs, title, subtitle):
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = title
    slide.placeholders[1].text = subtitle
    return slide

def add_content_slide(prs, title, bullet_points):
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = title

    text_frame = slide.placeholders[1].text_frame
    text_frame.clear()

    for i, point in enumerate(bullet_points):
        if i == 0:
            p = text_frame.paragraphs[0]
        else:
            p = text_frame.add_paragraph()

        if isinstance(point, tuple):
            text, level = point
            p.text = text
            p.level = level
        else:
            p.text = point
            p.level = 0

        p.font.size = Pt(18) if p.level == 0 else Pt(16)

    return slide

# Slide 1: Title
add_title_slide(prs,
    "Detecting Sentiment Cascades Across Stocks and Crypto Markets",
    "FIN 494: Financial Technology & Innovation\n\nCruz Flores")

# Slide 2: Project Topic
add_content_slide(prs, "Project Topic", [
    "Real-time Sentiment Cascade Detection Platform",
    ("Detecting rapid emotional waves in online investor discussions", 1),
    ("Analyzing how sentiment spreads across social media precedes market volatility", 1),
    "",
    "Core Innovation:",
    ("Hybrid sentiment analysis (VADER + FinBERT)", 1),
    ("Multi-signal cascade detection algorithm", 1),
    ("Adaptive thresholds for different asset classes", 1),
    "",
    "Deliverables:",
    ("Academic research paper analyzing predictive power", 1),
    ("Production-ready dashboard with real-time monitoring", 1),
    ("Comparative analysis: Stocks vs. Crypto sentiment behavior", 1)
])

# Slide 3: Background & Importance
add_content_slide(prs, "Background & Importance", [
    "The Rise of Retail Investors & Social Media",
    ("2021 GameStop phenomenon showed social media's market impact", 1),
    ("Reddit, Twitter, StockTwits now influential in price discovery", 1),
    ("Retail trading accounts for 20-25% of market volume", 1),
    "",
    "Why This Matters:",
    ("Traditional sentiment analysis too slow for modern markets", 1),
    ("Early detection of sentiment cascades = risk management tool", 1),
    ("Crypto vs. stocks provide unique behavioral comparison", 1),
    ("Academic gap: Limited research on multi-signal cascade detection", 1),
    "",
    "Potential Impact:",
    ("Help regulators detect market manipulation", 1),
    ("Enable traders to avoid pump-and-dump schemes", 1),
    ("Advance understanding of digital-era market dynamics", 1)
])

# Slide 4: Literature Review
add_content_slide(prs, "Prior Work in This Area", [
    "Academic Research:",
    ("Bollen et al. (2011): Twitter mood predicts stock market movements", 1),
    ("Sprenger et al. (2014): Relationship between tweets and trades", 1),
    ("Renault (2017): Intraday online investor sentiment analysis", 1),
    "",
    "Industry Applications:",
    ("Bloomberg Social Sentiment Index", 1),
    ("StockTwits sentiment indicators", 1),
    ("Crypto Fear & Greed Index", 1),
    "",
    "Gaps Our Project Addresses:",
    ("Most research focuses on aggregate sentiment, not velocity/cascades", 1),
    ("Limited cross-market comparison (stocks vs. crypto)", 1),
    ("Few studies on pump-and-dump pattern detection", 1),
    ("Lack of adaptive, asset-specific thresholds", 1)
])

# Slide 5: Research Paper Outline
add_content_slide(prs, "Research Paper Outline (15-20 pages)", [
    "1. Introduction & Literature Review (3-4 pages)",
    ("Problem statement and research questions", 1),
    ("Review of sentiment analysis in finance", 1),
    "",
    "2. Methodology (4-5 pages)",
    ("Data collection approach (Reddit, StockTwits, News)", 1),
    ("Hybrid sentiment analysis (VADER + FinBERT)", 1),
    ("Multi-signal cascade detection algorithm", 1),
    ("Adaptive baseline calculation", 1),
    "",
    "3. Data Analysis & Results (5-6 pages)",
    ("Cascade frequency and characteristics", 1),
    ("Predictive power analysis (sentiment → price volatility)", 1),
    ("Stocks vs. crypto comparison", 1),
    ("Pump-and-dump pattern detection results", 1),
    "",
    "4. Discussion & Conclusion (3-4 pages)",
    ("Interpretation of findings", 1),
    ("Limitations and future research", 1)
])

# Slide 6: Primary Research Questions
add_content_slide(prs, "Research Questions", [
    "Primary Question:",
    ("Can sudden shifts in online investor sentiment serve as early-warning", 1),
    ("indicators of abnormal market activity?", 1),
    ("Do they behave differently between stocks and crypto?", 1),
    "",
    "Secondary Questions:",
    ("When do sentiment cascades occur most frequently?", 1),
    ("How predictive are sentiment velocity spikes of price volatility?", 1),
    ("Do crypto cascades dissipate faster than stock cascades?", 1),
    ("Can we detect artificial hype (pump-and-dumps) using posting patterns?", 1),
    ("Does retail-institutional sentiment divergence predict reversals?", 1),
    "",
    "Key Hypotheses:",
    ("H1: Sentiment cascades precede volatility by 30-90 minutes", 1),
    ("H2: Crypto cascades have 2-3× faster velocity than stocks", 1),
    ("H3: Volume-weighted sentiment is more predictive than raw sentiment", 1)
])

# Slide 7: Data Sources
add_content_slide(prs, "Data Sources & Collection", [
    "Asset Coverage:",
    ("10 Stocks: GME, AMC, TSLA, NVDA, AAPL, MSFT, PLTR, BBBY, SPY, COIN", 1),
    ("5 Cryptocurrencies: BTC, ETH, DOGE, SHIB, SOL", 1),
    "",
    "Data Collection (Every 30 minutes):",
    ("Reddit API: r/wallstreetbets, r/stocks, r/cryptocurrency", 1),
    ("StockTwits API: Ticker-specific social sentiment", 1),
    ("NewsAPI: Financial news headlines", 1),
    ("Yahoo Finance (yfinance): Stock prices and volume", 1),
    ("CoinGecko API: Cryptocurrency prices and metrics", 1),
    "",
    "Data Volume Expected:",
    ("~10,000+ posts per week across all tickers", 1),
    ("3-4 weeks of continuous data collection", 1),
    ("All APIs offer free tiers sufficient for academic use", 1)
])

# Slide 8: Analysis Methodology
add_content_slide(prs, "Analysis Methodology", [
    "1. Sentiment Analysis (Hybrid Approach):",
    ("VADER: Fast, rule-based sentiment for social media posts", 1),
    ("FinBERT: Deep learning model for financial language", 1),
    ("Engagement-weighted scoring (upvotes + comments)", 1),
    "",
    "2. Cascade Detection Algorithm:",
    ("Calculate sentiment velocity (rate of change)", 1),
    ("Measure volume deltas (posting frequency spikes)", 1),
    ("Track engagement deltas (community activation)", 1),
    ("Multi-signal score: 40% sentiment + 35% volume + 25% engagement", 1),
    ("Adaptive thresholds based on each asset's volatility profile", 1),
    "",
    "3. Pattern Recognition:",
    ("Pump-and-dump detection: rapid positive spike + concentrated posting", 1),
    ("Retail vs. institutional divergence tracking", 1),
    "",
    "4. Statistical Analysis:",
    ("Correlation: sentiment velocity → price volatility (R² calculation)", 1),
    ("Lead time analysis: how far ahead do cascades predict moves?", 1),
    ("Comparative statistics: stocks vs. crypto behavior", 1)
])

# Slide 9: Questions We Can Answer
add_content_slide(prs, "Questions Answered by Our Analysis", [
    "Predictive Power:",
    ("Do sentiment cascades precede price volatility?", 1),
    ("What is the optimal lead time (1hr, 6hr, 24hr)?", 1),
    ("What correlation strength exists (R² target: >0.4 crypto, >0.3 stocks)?", 1),
    "",
    "Market Comparison:",
    ("Are crypto cascades faster/more frequent than stock cascades?", 1),
    ("Do different asset classes require different detection thresholds?", 1),
    ("Which market is more susceptible to social sentiment?", 1),
    "",
    "Pattern Detection:",
    ("Can we identify pump-and-dump schemes algorithmically?", 1),
    ("Does retail-institutional divergence predict reversals?", 1),
    ("Which tickers show strongest sentiment-price coupling?", 1),
    "",
    "Practical Applications:",
    ("Can retail traders use this as risk management tool?", 1),
    ("Could regulators detect market manipulation earlier?", 1)
])

# Slide 10: Technical Architecture
add_content_slide(prs, "Technical Implementation", [
    "Microservices Architecture:",
    ("Railway Platform: 5 microservices + PostgreSQL + Redis", 1),
    ("Vercel: Next.js dashboard frontend", 1),
    "",
    "Backend Services:",
    ("Data Collector: Python/FastAPI (Reddit, News, Market data)", 1),
    ("Sentiment Processor: VADER + FinBERT hybrid pipeline", 1),
    ("Cascade Detector: Multi-signal algorithm with adaptive thresholds", 1),
    ("Alert Engine: Real-time monitoring and notification system", 1),
    ("API Gateway: REST endpoints for dashboard", 1),
    "",
    "Dashboard Features:",
    ("Real-time sentiment vs. price charts", 1),
    ("Active cascade alerts (Yellow/Orange/Red levels)", 1),
    ("Stocks vs. crypto comparison view", 1),
    ("Historical playback and data export for research", 1),
    "",
    ("Total Project Cost: $25-50 for 4 weeks (using free tiers + Railway credits)", 0)
])

# Slide 11: Timeline
add_content_slide(prs, "Project Timeline (4 Weeks)", [
    "Week 1: Foundation & Data Collection (Nov 4-10)",
    ("Deploy Railway infrastructure and PostgreSQL database", 1),
    ("Build data collectors for Reddit, StockTwits, NewsAPI", 1),
    ("Begin collecting 7 days of baseline data", 1),
    ("Milestone: Data Pipeline Operational (10k+ posts collected)", 1),
    "",
    "Week 2: Sentiment & Cascade Detection (Nov 11-17)",
    ("Implement VADER + FinBERT hybrid sentiment analysis", 1),
    ("Build cascade detection algorithm with adaptive thresholds", 1),
    ("Generate first week of cascade events", 1),
    ("Milestone: Cascade Detection Live (20+ cascades detected)", 1),
    "",
    "Week 3: Dashboard & Analysis (Nov 18-24)",
    ("Deploy Vercel dashboard with core features", 1),
    ("Implement pump-and-dump pattern detection", 1),
    ("Conduct initial research analysis (correlations, statistics)", 1),
    ("Milestone: Dashboard Demo-Ready", 1)
])

# Slide 12: Timeline Continued
add_content_slide(prs, "Project Timeline - Completion", [
    "Week 4: Polish & Research Paper (Nov 25 - Dec 1)",
    ("Complete data export functionality", 1),
    ("Draft 15-20 page research paper", 1),
    ("Create data visualizations for paper", 1),
    ("Prepare final presentation and demo video", 1),
    ("Final testing and bug fixes", 1),
    ("Milestone: Project Complete - Ready for submission", 1),
    "",
    "Current Status:",
    ("Week 1 in progress: Infrastructure setup underway", 1),
    ("Data collection starting this week", 1),
    ("On track for 3+ weeks of data before final presentation", 1),
    "",
    "Success Metrics:",
    ("50+ cascade events detected across all tickers", 1),
    (">60% cascade detection precision", 1),
    ("Measurable difference in stock vs. crypto behavior", 1)
])

# Slide 13: Expected Outcomes
add_content_slide(prs, "Expected Outcomes & Contributions", [
    "Academic Contributions:",
    ("Novel multi-signal cascade detection methodology", 1),
    ("First systematic comparison of stock vs. crypto sentiment cascades", 1),
    ("Adaptive threshold algorithm for different asset volatility profiles", 1),
    ("Pump-and-dump detection framework", 1),
    "",
    "Practical Value:",
    ("Working dashboard for real-time sentiment monitoring", 1),
    ("Exportable dataset for future research", 1),
    ("Open-source codebase for reproducibility", 1),
    "",
    "Potential Impact:",
    ("Risk management tool for retail traders", 1),
    ("Market surveillance application for regulators", 1),
    ("Foundation for future FinTech startup", 1),
    "",
    "Publication Potential:",
    ("Suitable for finance/data science journals", 1),
    ("Presentation at student research conferences", 1)
])

# Slide 14: Risks & Mitigations
add_content_slide(prs, "Risks & Mitigation Strategies", [
    "Technical Risks:",
    ("Reddit API rate limits → Use PRAW library with exponential backoff", 1),
    ("FinBERT processing slow → Batch processing, fallback to VADER only", 1),
    ("Railway costs exceed budget → Monitor daily, reduce frequency if needed", 1),
    "",
    "Research Risks:",
    ("False positive cascades → Document honestly, tune thresholds with data", 1),
    ("Low correlation results → Still academically valuable (null results matter)", 1),
    ("Insufficient data quality → Validation checks, engagement thresholds", 1),
    "",
    "Timeline Risks:",
    ("Not enough time for features → Focus on MUST-HAVEs first (MoSCoW)", 1),
    ("Paper + basic dashboard > fancy features alone", 1),
    "",
    "Overall Risk Level: MEDIUM",
    ("Well-scoped project with clear priorities", 1),
    ("Realistic timeline with buffer built in", 1)
])

# Slide 15: Q&A
slide = prs.slides.add_slide(prs.slide_layouts[1])
slide.shapes.title.text = "Questions & Discussion"

text_frame = slide.placeholders[1].text_frame
text_frame.clear()

questions = [
    "Discussion Topics:",
    ("What additional data sources would strengthen this analysis?", 1),
    ("Are there specific tickers or market events we should focus on?", 1),
    ("What academic journals would be most interested in these findings?", 1),
    ("Should we consider adding Twitter/X data (requires paid API)?", 1),
    "",
    "",
    "Contact Information:",
    ("Student: Cruz Flores", 1),
    ("Course: FIN 494 - Financial Technology & Innovation", 1),
    ("Project Repository: [Will be provided]", 1),
    ("Dashboard URL: [Will be deployed to Vercel]", 1)
]

for i, point in enumerate(questions):
    if i == 0:
        p = text_frame.paragraphs[0]
    else:
        p = text_frame.add_paragraph()

    if isinstance(point, tuple):
        text, level = point
        p.text = text
        p.level = level
    else:
        p.text = point
        p.level = 0

    p.font.size = Pt(18) if p.level == 0 else Pt(16)

# Save presentation
output_path = "/Users/cruzflores/Desktop/FIN 494/Sentiment_Cascade_Detection_Presentation.pptx"
prs.save(output_path)
print(f"Presentation created successfully: {output_path}")
print(f"Total slides: {len(prs.slides)}")
